from pptx import Presentation

# Lade die Präsentation
prs = Presentation('Template.pptx')

# Zugriff auf alle verfügbaren Masterlayouts
for slide_layout in prs.slide_layouts:
    print(slide_layout.name)
    # Du kannst hier weitere Eigenschaften jedes Layouts untersuchen


# Zugriff auf den Präsentationsmaster und die Masterlayouts
master = prs.slide_master
for shape in master.shapes:
    if shape.has_text_frame:
        print(shape.text)
        # Untersuche Schriftstil, Größe und Farbe