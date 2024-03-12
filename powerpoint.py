from pptx import Presentation
from pptx.util import Inches
import pandas as pd

def create_presentation(dataframes):
    prs = Presentation()

    for heat_name, df in dataframes.items():
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # 5 für Titel und Inhalt
        title = slide.shapes.title
        title.text = heat_name

        # Tabelle hinzufügen
        rows, cols = df.shape[0] + 1, df.shape[1]
        table = slide.shapes.add_table(rows, cols, Inches(2), Inches(1.5), Inches(6), Inches(0.8 * rows)).table

        # Spaltenüberschriften
        for col_index, col_name in enumerate(df.columns):
            table.cell(0, col_index).text = col_name

        # Tabelledaten
        for row_index, row in df.iterrows():
            for col_index, item in enumerate(row):
                table.cell(row_index + 1, col_index).text = str(item)

    # Speichere die Präsentation
    prs.save('Automatisierte_Präsentation.pptx')

# Beispiel: DataFrames
dataframes = {
    'Heat 1': pd.DataFrame({
        'Name': ['Athlet 1', 'Athlet 2'],
        'Zeit': ['10.23s', '10.36s']
    }),
    # Weitere DataFrames hinzufügen
}

create_presentation(dataframes)